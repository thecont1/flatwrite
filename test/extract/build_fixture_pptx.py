"""
Build a minimal PPTX fixture for testing the powerpoint_notes rule.

Run from repo root:
    uv run --directory services/extract python test/extract/build_fixture_pptx.py

The fixture has:
  - 3 slides
  - Slide 1: has both bullets and speaker notes
  - Slide 2: has bullets only, no notes (tests the "fall back to full
    content" branch of the rule)
  - Slide 3: has notes only, no body bullets
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

OUT_PATH = Path(__file__).resolve().parent / "fixtures" / "sample.pptx"


def build() -> Path:
    prs = Presentation()
    # Slide 1: bullets + notes
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Welcome"
    body = slide1.placeholders[1]
    body.text_frame.text = "First bullet"
    # second paragraph
    p = body.text_frame.add_paragraph()
    p.text = "Second bullet"
    notes_tf = slide1.notes_slide.notes_text_frame
    notes_tf.text = "Speaker: greet the audience warmly and pause for laughs."
    notes_tf.add_paragraph().text = "Next: introduce the agenda."

    # Slide 2: bullets only, no notes
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1]
    body2.text_frame.text = "Item one"
    p = body2.text_frame.add_paragraph()
    p.text = "Item two"

    # Slide 3: notes only
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Closing"
    notes_tf3 = slide3.notes_slide.notes_text_frame
    notes_tf3.text = "Wrap up with the call to action and thank-yous."

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    p = build()
    print(f"wrote {p} ({p.stat().st_size} bytes)")
